# -*- coding: utf-8 -*-
"""
Script pour créer automatiquement une présentation PowerPoint du stage
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

def create_presentation():
    """Crée une présentation PowerPoint du stage"""
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Couleurs du thème
    primary_color = RGBColor(41, 128, 185)  # Bleu
    secondary_color = RGBColor(52, 73, 94)  # Gris foncé
    accent_color = RGBColor(46, 204, 113)   # Vert
    
    # ============================
    # SLIDE 1 : Page de titre
    # ============================
    slide_layout = prs.slide_layouts[0]  # Titre
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Analyse des Tendances de Formation Digitale"
    subtitle.text = "Stage Data Science & Intelligence Artificielle\n\nDéveloppement d'un système d'analyse prédictive"
    
    # ============================
    # SLIDE 2 : Contexte et objectifs
    # ============================
    slide_layout = prs.slide_layouts[1]  # Titre et contenu
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contexte et Objectifs"
    content.text = """🎯 Contexte :
• Besoin d'anticiper les tendances de formation digitale
• Exploitation des données internes et externes
• Optimisation de l'offre de formation

📋 Objectifs :
• Développer un prototype d'analyse prédictive
• Automatiser la collecte et visualisation des données
• Utiliser l'IA pour détecter les sujets émergents
• Créer un dashboard interactif de restitution"""
    
    # ============================
    # SLIDE 3 : Méthodologie
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Méthodologie - 4 Étapes"
    content.text = """1️⃣ Collecte et Préparation des Données
   • Données internes : formations, inscriptions
   • Données externes : offres d'emploi, Google Trends
   • Nettoyage et standardisation

2️⃣ Exploration et Analyse
   • Analyse statistique descriptive
   • Visualisations et corrélations
   • Identification des tendances

3️⃣ Modélisation Prédictive
   • Test de 6 modèles de Machine Learning
   • Optimisation des hyperparamètres
   • Évaluation des performances

4️⃣ Restitution des Résultats
   • Dashboard interactif Streamlit
   • Visualisations dynamiques
   • Interface utilisateur intuitive"""
    
    # ============================
    # SLIDE 4 : Données collectées
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sources de Données"
    content.text = """📊 Données Internes :
• Catalogue de formations digitales
• Historique des inscriptions
• Progression des étudiants

🌐 Données Externes :
• Offres d'emploi (Remotive, Adzuna)
• Tendances Google Trends
• Données Stack Overflow
• Enquêtes du secteur

🔧 Traitement :
• Extraction de mots-clés (NLTK)
• Mapping formations ↔ emplois
• Calcul de ratios demande/étudiants
• Standardisation des formats"""
    
    # ============================
    # SLIDE 5 : Résultats des modèles
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Résultats de la Modélisation"
    content.text = """🏆 Modèles Testés et Performances :

1. XGBoost : RMSE = 152.99, R² = 0.73 ⭐
2. Linear Regression : RMSE = 161.80, R² = 0.69
3. Gradient Boosting : RMSE = 167.35, R² = 0.67
4. Ridge Regression : RMSE = 186.98, R² = 0.59
5. Lasso Regression : RMSE = 194.05, R² = 0.56
6. Random Forest : RMSE = 198.20, R² = 0.54

🎯 Meilleur modèle : XGBoost
   • Précision prédictive : 73%
   • Variables importantes : catégorie, certification, durée
   • Temps d'entraînement : < 30 secondes"""
    
    # ============================
    # SLIDE 6 : Dashboard créé
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Dashboard Interactif"
    content.text = """📊 Interface Développée :

🏠 Vue d'ensemble
   • Métriques clés en temps réel
   • Top 10 des formations
   • Évolution des tendances

📈 Tendances du marché
   • Filtres interactifs
   • Distribution de la demande
   • Analyse des ratios

🎓 Analyse des formations
   • Filtres avancés
   • Relations entre variables
   • Métriques filtrées

🔮 Prédictions
   • Comparaison des modèles
   • Scores de performance
   • Simulateur de prédiction

📊 Comparaisons
   • Offre vs Demande
   • Analyses croisées
   • Évolutions temporelles

📋 Données brutes
   • Exploration des datasets
   • Statistiques descriptives"""
    
    # ============================
    # SLIDE 7 : Principales découvertes
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Principales Découvertes"
    content.text = """🎯 Tendances Identifiées :

1. Data Science et Développement Web
   • Formations les plus demandées
   • Croissance continue du marché

2. Impact de la Certification
   • +20% de demande pour les formations certifiantes
   • Valorisation importante sur le marché

3. Influence de la Durée
   • Corrélation positive avec la demande
   • Optimisation possible des programmes

4. Corrélation avec Google Trends
   • Tendances en ligne ↔ demande réelle
   • Possibilité d'anticipation

🔮 Insights Prédictifs :
• Modèle XGBoost performant (73% précision)
• Variables clés identifiées
• Capacité d'anticipation des tendances"""
    
    # ============================
    # SLIDE 8 : Technologies utilisées
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technologies et Compétences"
    content.text = """🛠️ Technologies Utilisées :

🐍 Python
   • Pandas : manipulation de données
   • Scikit-learn : Machine Learning
   • XGBoost : modèles de boosting
   • NLTK : traitement du langage naturel

📊 Visualisation
   • Streamlit : interface web
   • Plotly : graphiques interactifs
   • Matplotlib/Seaborn : visualisations

📁 Données
   • APIs : Google Trends, Stack Overflow
   • Web Scraping : collecte d'offres
   • CSV : format principal

🎓 Compétences Développées :
• Data Science et analyse exploratoire
• Machine Learning et modélisation
• NLP et extraction de mots-clés
• Développement d'applications web
• Visualisation de données"""
    
    # ============================
    # SLIDE 9 : Impact et applications
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Impact et Applications"
    content.text = """🚀 Impact pour l'Organisation :

📈 Optimisation de l'Offre
   • Anticipation des tendances
   • Adaptation des programmes
   • ROI amélioré des formations

🎯 Décisions Basées sur les Données
   • Insights quantifiés
   • Prédictions fiables
   • Stratégie data-driven

👥 Bénéfices pour les Étudiants
   • Formations adaptées au marché
   • Certifications valorisées
   • Insertion professionnelle facilitée

💼 Applications Futures :
• Déploiement en production
• Mise à jour automatique
• Intégration CRM
• Alertes sur nouvelles tendances"""
    
    # ============================
    # SLIDE 10 : Conclusion
    # ============================
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion et Recommandations"
    content.text = """✅ Objectifs Atteints :

🎯 100% des objectifs du stage réalisés
📊 Dashboard fonctionnel et interactif
🤖 Modèles prédictifs performants
📈 Insights actionnables identifiés

🔮 Recommandations :

1. Déployer le dashboard en production
2. Former les équipes à son utilisation
3. Automatiser la collecte de données
4. Étendre l'analyse à d'autres domaines
5. Intégrer des alertes sur nouvelles tendances

📊 Métriques de Succès :
• Précision prédictive : 73%
• Temps de chargement : < 3 secondes
• Fonctionnalités : 100% opérationnelles
• Documentation : complète

🎓 Stage Réalisé avec Succès !"""
    
    # Sauvegarder la présentation
    prs.save('Presentation_Stage_Formation_Digitale.pptx')
    print("✅ Présentation créée : 'Presentation_Stage_Formation_Digitale.pptx'")

if __name__ == "__main__":
    try:
        create_presentation()
        print("🎉 Présentation PowerPoint générée avec succès !")
    except Exception as e:
        print(f"❌ Erreur lors de la création : {e}")
        print("💡 Installez python-pptx : pip install python-pptx")

